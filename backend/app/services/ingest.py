"""
Document ingestion: extract text, split, add to vector store
"""
import os
import logging
from typing import List
from uuid import uuid4
from app.core.config import load_config
from app.services.vector_store import add_documents
import datetime

logger = logging.getLogger(__name__)


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    text_parts = []
    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader

            reader = PdfReader(path)
            for p in reader.pages:
                text_parts.append(p.extract_text() or "")
        elif ext == ".docx":
            from docx import Document

            doc = Document(path)
            for p in doc.paragraphs:
                text_parts.append(p.text)
        elif ext in (".txt", ".md"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text_parts.append(f.read())
        elif ext in (".pptx",):
            from pptx import Presentation

            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
        elif ext in (".xlsx", ".xls"):
            # 使用 MarkItDown 进行 Excel 转 Markdown
            text_parts.append(extract_excel_as_markdown(path))
        else:
            # fallback: try to read as text
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text_parts.append(f.read())
    except Exception as e:
        logger.error(f"提取文本失败 ({path}): {e}")
    return "\n".join(text_parts)


def extract_excel_as_markdown(excel_path: str) -> str:
    """
    使用 MarkItDown 将 Excel 文件转换为 Markdown 格式

    Args:
        excel_path: Excel 文件路径

    Returns:
        Markdown 格式的字符串
    """
    try:
        # 首先尝试使用 MarkItDown
        try:
            from markitdown import MarkItDown

            logger.info(f"使用 MarkItDown 转换 Excel 文件: {excel_path}")
            md = MarkItDown()
            result = md.convert(excel_path)
            return result.text_content

        except ImportError:
            logger.warning("MarkItDown 未安装，尝试使用备选方案")
        except Exception as e:
            logger.warning(f"MarkItDown 转换失败: {e}，尝试备选方案")

        # 备选方案：使用 openpyxl
        try:
            import openpyxl

            logger.info(f"使用 openpyxl 转换 Excel 文件: {excel_path}")
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            markdown_lines = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                markdown_lines.append(f"## {sheet_name}\n")

                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    if all(cell is None for cell in row):
                        continue

                    cells = []
                    for cell in row:
                        if cell is None:
                            cells.append("")
                        elif isinstance(cell, (int, float)):
                            cells.append(str(cell))
                        else:
                            cell_str = str(cell)
                            cell_str = cell_str.replace('|', '\\|')
                            cells.append(cell_str.strip())

                    markdown_lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 1:
                        markdown_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

                markdown_lines.append("")
            return "\n".join(markdown_lines)

        except ImportError:
            logger.error("openpyxl 未安装")
            raise Exception("请安装 openpyxl 或 markitdown 库以支持 Excel 文件处理")
        except Exception as e:
            logger.error(f"openpyxl 转换失败: {e}")
            raise

    except Exception as e:
        logger.error(f"Excel 文件处理失败: {e}")
        raise


def split_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    if not text:
        return []

    sections = _split_by_sections(text)
    if sections and len(sections) > 1:
        chunks = _merge_sections_to_chunks(sections, chunk_size, chunk_overlap)
        if chunks:
            return chunks

    chunks = []
    step = chunk_size - chunk_overlap
    if step <= 0:
        step = chunk_size
    for i in range(0, len(text), step):
        chunk = text[i : i + chunk_size]
        chunks.append(chunk)
    return chunks


def _split_by_sections(text: str) -> List[dict]:
    """
    按论文章节结构拆分文本，保留章节标题信息

    Returns:
        章节列表，每个元素包含 title 和 content
    """
    import re

    section_patterns = [
        r'^\s*(\d+\.?\s+[A-Z][a-zA-Z\s,:\-]+)',
        r'^\s*(Abstract|Introduction|Related Work|Method|Methodology|Experiment|Results|Discussion|Conclusion|References|Acknowledgment)',
        r'^\s*(摘要|引言等相关工作|方法|实验|结果|讨论|结论|参考文献|致谢)',
        r'^\s*(\d+\.\d+\s+.+)',
    ]

    combined_pattern = '|'.join(f'({p})' for p in section_patterns)

    lines = text.split('\n')
    sections = []
    current_title = ""
    current_content = []

    for line in lines:
        is_section_header = False
        for pattern in section_patterns:
            if re.match(pattern, line.strip()) and len(line.strip()) < 100:
                is_section_header = True
                break

        if is_section_header and current_content:
            sections.append({
                "title": current_title,
                "content": "\n".join(current_content).strip()
            })
            current_title = line.strip()
            current_content = []
        elif is_section_header:
            current_title = line.strip()
        else:
            current_content.append(line)

    if current_content:
        sections.append({
            "title": current_title,
            "content": "\n".join(current_content).strip()
        })

    return sections


def _merge_sections_to_chunks(sections: List[dict], chunk_size: int, chunk_overlap: int) -> List[str]:
    """
    将章节合并为合适大小的chunks，保留章节标题上下文
    """
    chunks = []
    current_chunk = ""

    for section in sections:
        title = section["title"]
        content = section["content"]

        section_text = f"[{title}]\n{content}" if title else content

        if len(current_chunk) + len(section_text) + 2 <= chunk_size:
            if current_chunk:
                current_chunk += "\n\n" + section_text
            else:
                current_chunk = section_text
        else:
            if current_chunk:
                chunks.append(current_chunk)

            if len(section_text) <= chunk_size:
                current_chunk = section_text
            else:
                step = chunk_size - chunk_overlap
                if step <= 0:
                    step = chunk_size
                for i in range(0, len(section_text), step):
                    chunk_part = section_text[i : i + chunk_size]
                    if chunk_part.strip():
                        chunks.append(chunk_part)
                current_chunk = ""

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def generate_and_save_summary(
    file_path: str,
    filename: str,
    file_text: str,
    provider: str = "openai"
) -> bool:
    """
    生成并保存文件摘要

    Args:
        file_path: 文件路径
        filename: 文件名
        file_text: 文件文本内容
        provider: LLM提供商

    Returns:
        是否成功
    """
    cfg = load_config()
    summary_config = cfg.get("summary_search", {})

    if not summary_config.get("auto_generate_summary", True):
        logger.info("自动生成摘要功能已关闭")
        return False

    try:
        from app.services.document_summary import generate_document_summary
        from app.services.summary_store import save_document_summary

        logger.info("Generating summary for: %s", filename)

        summary = generate_document_summary(file_text, filename, provider)
        summary.file_path = os.path.relpath(file_path, os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))

        # 保存摘要
        success = save_document_summary(summary)

        if success:
            logger.info("Summary saved: %s (score: %.2f)", filename, summary.quality_score)
        else:
            logger.error("Summary save failed: %s", filename)

        return success

    except Exception as e:
        logger.error("Summary generation error: %s", e, exc_info=True)
        return False


def ingest_file(path: str, provider: str = "openai", generate_summary: bool = True) -> dict:
    cfg = load_config()
    text = extract_text(path)
    filename = os.path.basename(path)

    result = {
        "filename": filename,
        "chunks_count": 0,
        "parent_count": 0,
        "summary_generated": False,
        "success": False
    }

    try:
        from app.services.vector_store import delete_documents_by_filename
        deleted = delete_documents_by_filename(filename)
        if deleted > 0:
            logger.debug("Deleted %d old chunks for %s before ingest", deleted, filename)
    except Exception as e:
        logger.warning(f"摄入前去重删除失败: {e}")

    if generate_summary:
        summary_success = generate_and_save_summary(path, filename, text, provider)
        result["summary_generated"] = summary_success

    pdr_config = cfg.get("parent_document_retrieval", {})
    use_parent_mode = pdr_config.get("enabled", False)

    if use_parent_mode and (filename.endswith('.md') or filename.endswith('.txt')):
        from app.services.parent_document import generate_parent_documents, generate_child_chunks
        from app.services.parent_store import ParentDocumentStore

        parent_chunk_size = pdr_config.get("parent_max_chars", 8000)
        child_chunk_size = pdr_config.get("child_chunk_size", 300)
        child_chunk_overlap = pdr_config.get("child_chunk_overlap", 60)

        parents = generate_parent_documents(
            text, filename,
            min_parent_chars=pdr_config.get("parent_min_chars", 300),
            max_parent_chars=parent_chunk_size
        )
        result["parent_count"] = len(parents)

        all_children = []
        for parent in parents:
            children = generate_child_chunks(
                parent,
                child_chunk_size=child_chunk_size,
                child_chunk_overlap=child_chunk_overlap
            )
            all_children.extend(children)

        parent_store = ParentDocumentStore()
        parent_store.delete_by_filename(filename)
        parent_store.save_parents(parents)

        child_ids = [f"{c['parent_id']}_{c['chunk_index']}" for c in all_children]
        child_texts = [c['text'] for c in all_children]
        child_metas = [
            {
                "source": c['filename'],
                "chunk_index": c['chunk_index'],
                "parent_id": c['parent_id'],
                "section_title": c['section_title'],
                "chunk_type": "child",
                "ingest_time": datetime.datetime.utcnow().isoformat(),
                "file_type": os.path.splitext(path)[1].lower(),
            }
            for c in all_children
        ]

        if child_texts:
            add_documents(ids=child_ids, documents=child_texts, metadatas=child_metas, provider=provider)
            result["chunks_count"] = len(child_texts)
            result["success"] = True
            logger.info(
                "Ingest complete (parent mode): %s → %d parents + %d children",
                filename, len(parents), len(child_texts)
            )

    else:
        chunks = split_text(text, cfg.get("chunk_size", 1500), cfg.get("chunk_overlap", 225))
        metas = []
        ids = []
        for i, c in enumerate(chunks):
            metas.append({
                "source": filename,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "chunk_type": "legacy",
                "ingest_time": datetime.datetime.utcnow().isoformat(),
                "file_type": os.path.splitext(path)[1].lower(),
            })
            ids.append(str(uuid4()))
        if chunks:
            add_documents(ids=ids, documents=chunks, metadatas=metas, provider=provider)
            result["chunks_count"] = len(chunks)
            result["success"] = True
            logger.info("Ingest complete (legacy mode): %s → %d chunks", filename, len(chunks))

    try:
        from app.services.bm25_search import refresh_bm25_index
        refresh_bm25_index()
        logger.debug("BM25 index refreshed after ingest")
    except Exception as e:
        logger.warning("BM25 index refresh failed after ingest: %s", e)

    return result
